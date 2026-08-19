"""End-to-end and unit tests for the custom track.

Run: pytest -q
"""
from __future__ import annotations

import asyncio

import pytest


def _run(coro):
    return asyncio.run(coro)

from docextract.custom.extractor import Extractor
from docextract.custom.reviewer import review, apply_decision, decision_row
from docextract.custom.storage import DeltaLakeStorage
from docextract.custom.validator import reconcile, gate_for_review
from docextract.shared.dictionary import canonical_metric
from docextract.shared.gates import (
    SignOff, SecondLineValidation, Stage,
    gate_external_release, gate_production_deployment,
)
from docextract.shared.cost_tracker import CostTracker
from docextract.shared.llm_client import LLMResponse
from docextract.shared.llm_router import QueryRouter, is_safe_select
from docextract.shared.routing import decide_start_tier
from docextract.shared.schema import Claim, Element, Modality, parse_numeric
from tests.stubs import StubLLMClient, FakeSearchStore


def _element(content: str, entity_ref="FIRM123", team="capital"):
    return Element(element_id="e1", modality=Modality.TEXT, content=content,
                   source_document="/x/doc.pdf", team=team,
                   report_type="primary_report", entity_ref=entity_ref)


# --- extraction alignment -------------------------------------------------

def test_extract_populates_all_structural_fields():
    client = StubLLMClient(default_claims=[{
        "field_name": "Primary ratio", "value": "14.2", "unit": "%",
        "reporting_basis": "consolidated", "netting": "net of deductions",
        "scale": "ratio", "as_at_date": "2025-12-31", "confidence": 0.95,
    }])
    res = _run(Extractor(client=client).extract([_element("has the primary metric")]))
    assert len(res.claims) == 1
    c = res.claims[0]
    assert c.canonical_metric == "primary_ratio"
    assert c.unit == "%" and c.reporting_basis == "consolidated"
    assert c.netting == "net of deductions" and c.scale == "ratio"
    assert c.as_at_date == "2025-12-31"
    assert c.entity_ref == "FIRM123"           # propagated from element
    assert c.model_used == "tier1"


def test_low_confidence_escalates_to_tier2():
    client = StubLLMClient(extract_script={
        "trigger": [{"field_name": "coverage", "value": "120", "confidence": 0.4}],
    }, default_claims=[])
    # tier1 returns conf 0.4 -> escalate; tier2 script keyed on same content.
    res = _run(Extractor(client=client).extract([_element("trigger coverage")]))
    assert res.claims[0].model_used == "tier2"


def test_malformed_row_is_quarantined_not_fatal():
    client = StubLLMClient(default_claims=[
        {"value": "10", "confidence": 0.9},          # missing field_name
        {"field_name": "coverage", "value": "120", "confidence": 0.9},  # good
    ])
    res = _run(Extractor(client=client).extract([_element("mix")]))
    assert len(res.claims) == 1
    assert len(res.quarantine) == 1


# --- content-aware routing ------------------------------------------------

def _el(content="text here", modality=Modality.TEXT,
        report_type="primary_report"):
    return Element(element_id="e", modality=modality, content=content,
                   source_document="/x/doc.pdf", team="capital",
                   report_type=report_type, entity_ref="FIRM123")


class _RecordingClient:
    """Records the tier of every extract call; always high confidence so the
    reactive escalation would NOT fire on its own."""

    def __init__(self):
        self.calls: list[str] = []

    async def extract(self, model_key, prompt, content, tool_schema, image_base64=None):
        self.calls.append(model_key)
        return LLMResponse(text="", tool_arguments={"claims": [
            {"field_name": "Primary ratio", "value": "14.2", "confidence": 0.95}]},
            input_tokens=100, output_tokens=50)


def test_route_decision_by_modality_and_length():
    assert decide_start_tier(_el(modality=Modality.TABLE)).tier == "tier2"
    assert decide_start_tier(_el(modality=Modality.TEXT)).tier == "tier1"
    assert decide_start_tier(_el(content="x" * 7000)).tier == "tier2"
    # Empty/flagged elements gain nothing from Tier 2.
    empty = decide_start_tier(_el(content="   ", modality=Modality.IMAGE))
    assert empty.tier == "tier1"
    # Decisions are auditable.
    d = decide_start_tier(_el(modality=Modality.TABLE))
    assert d.content_aware is True and "modality" in d.reason
    assert decide_start_tier(_el()).content_aware is False


def test_hard_content_skips_tier1_pass():
    """A table starts on Tier 2 directly — no wasted Tier 1 call."""
    rc = _RecordingClient()
    _run(Extractor(client=rc).extract([_el(modality=Modality.TABLE)]))
    assert rc.calls == ["tier2"]


def test_easy_content_stays_tier1_when_confident():
    rc = _RecordingClient()
    _run(Extractor(client=rc).extract([_el(modality=Modality.TEXT)]))
    assert rc.calls == ["tier1"]


def test_reactive_cascade_intact_for_low_confidence_text():
    """Text that comes back low-confidence still escalates Tier 1 -> Tier 2."""
    class _LowConf:
        def __init__(self):
            self.calls = []

        async def extract(self, model_key, prompt, content, tool_schema, image_base64=None):
            self.calls.append(model_key)
            conf = 0.4 if model_key == "tier1" else 0.9
            return LLMResponse(text="", tool_arguments={"claims": [
                {"field_name": "coverage", "value": "120", "confidence": conf}]},
                input_tokens=100, output_tokens=50)

    lc = _LowConf()
    _run(Extractor(client=lc).extract([_el(modality=Modality.TEXT)]))
    assert lc.calls == ["tier1", "tier2"]


def test_routing_recorded_on_result():
    rc = _RecordingClient()
    res = _run(Extractor(client=rc).extract([_el(modality=Modality.TABLE)]))
    assert res.routing["e"].tier == "tier2"


def test_cost_projection_accounts_for_hard_share():
    # hard_share=0 reduces to the reactive-only model.
    base = CostTracker.project(1000, 0.01, 0.20, 0.05)
    assert base == round(1000 * 0.01 + 1000 * 0.20 * 0.05, 2)
    # 30% routed straight to Tier 2 pay Tier 2 only; the rest cascade.
    mixed = CostTracker.project(1000, 0.01, 0.20, 0.05, hard_share=0.30)
    expected = round(1000 * 0.30 * 0.05 + 700 * 0.01 + 700 * 0.20 * 0.05, 2)
    assert mixed == expected


# --- OCR fallback for scanned pages ---------------------------------------

from docextract.custom import preprocessor as _pre
from docextract.custom.ai_parse_client import (
    ParsedBlock, _split_parsed, _is_table_sep)


class _FakePage:
    def to_image(self, resolution=200):
        class _Img:
            def save(self, buf, format):
                buf.write(b"PNGDATA")
        return _Img()


class _FakeOCR:
    def __init__(self, available=True, blocks=None, raise_exc=False):
        self.available = available
        self._blocks = blocks or []
        self._raise = raise_exc

    def parse_page_image(self, image_bytes, mime="image/png"):
        if self._raise:
            raise RuntimeError("endpoint down")
        return self._blocks


_OCR_ARGS = (_FakePage(), 3, "/x/scan.pdf", "capital", "primary_report",
             "FIRM1")


def test_split_parsed_detects_tables():
    blocks = _split_parsed("Heading\n\n| a | b |\n|---|---|\n| 1 | 2 |")
    assert len(blocks) == 2
    assert blocks[0].is_table is False and blocks[1].is_table is True
    assert _is_table_sep("|---|---|") and not _is_table_sep("plain line")


def test_ocr_unavailable_keeps_flagged_empty():
    els = _pre._ocr_page(*_OCR_ARGS, _FakeOCR(available=False))
    assert len(els) == 1 and els[0].modality == Modality.IMAGE
    assert els[0].content == ""


def test_ocr_empty_result_keeps_flagged_empty():
    els = _pre._ocr_page(*_OCR_ARGS, _FakeOCR(blocks=[]))
    assert len(els) == 1 and els[0].content == ""


def test_ocr_exception_keeps_flagged_empty():
    """A transport/parse failure must not lose the page."""
    els = _pre._ocr_page(*_OCR_ARGS, _FakeOCR(raise_exc=True))
    assert len(els) == 1 and els[0].modality == Modality.IMAGE


def test_ocr_success_emits_real_elements():
    ocr = _FakeOCR(blocks=[
        ParsedBlock("Recovered text", is_table=False),
        ParsedBlock("| a |\n|---|\n| 1 |", is_table=True)])
    els = _pre._ocr_page(*_OCR_ARGS, ocr)
    assert len(els) == 2
    assert els[0].modality == Modality.TEXT and els[0].content == "Recovered text"
    assert els[1].modality == Modality.TABLE
    assert "ocr" in els[0].element_id            # provenance visible in id
    assert els[0].entity_ref == "FIRM1" and els[0].page == 3
    assert all(e.content for e in els)           # no flagged-empty on success


# --- dictionary -----------------------------------------------------------

def test_canonical_metric_collapses_whitespace():
    assert canonical_metric("primary  ratio")[0] == "primary_ratio"
    name, mapped = canonical_metric("some novel metric")
    assert mapped is False and name == "some_novel_metric"


# --- value parsing / reconciliation --------------------------------------

def test_parse_numeric_handles_formatting():
    assert parse_numeric("14.2%") == 14.2
    assert parse_numeric("1,234") == 1234.0
    assert parse_numeric("£5m") == 5.0
    assert parse_numeric("n/a") is None


def test_reconcile_respects_structural_signature():
    # Same metric, different basis -> NOT compared -> no magnitude conflict.
    claims = [
        Claim(field_name="the primary metric", canonical_metric="primary_ratio", value="14",
              reporting_basis="consolidated", source_element_id="e1"),
        Claim(field_name="the primary metric", canonical_metric="primary_ratio", value="40",
              reporting_basis="solo", source_element_id="e2"),
    ]
    assert reconcile(claims) == []
    # Same signature, >2x apart -> conflict.
    claims2 = [
        Claim(field_name="the primary metric", canonical_metric="primary_ratio", value="14",
              reporting_basis="consolidated", source_element_id="e1"),
        Claim(field_name="the primary metric", canonical_metric="primary_ratio", value="40",
              reporting_basis="consolidated", source_element_id="e2"),
    ]
    assert len(reconcile(claims2)) == 1


def test_gate_for_review_flags_low_conf_and_conflict():
    claims = [
        Claim(field_name="A", canonical_metric="a", value="1",
              confidence=0.5, source_element_id="e1"),   # low conf
        Claim(field_name="B", canonical_metric="b", value="2",
              confidence=0.99, source_element_id="e2"),  # clean
    ]
    gated = gate_for_review(claims, conflicts=[])
    assert len(gated["needs_review"]) == 1 and len(gated["clean"]) == 1


# --- SQL guard ------------------------------------------------------------

@pytest.mark.parametrize("sql,ok", [
    ("SELECT * FROM gold_metrics", True),
    ("SELECT delete_flag FROM created_returns", True),   # no false positive
    ("SELECT * FROM t; DROP TABLE t", False),
    ("DELETE FROM gold_metrics", False),
    ("UPDATE gold_metrics SET value='0'", False),
])
def test_is_safe_select(sql, ok):
    assert is_safe_select(sql) is ok


def test_router_generates_sql_when_absent():
    client = StubLLMClient(
        complete_script={"the primary metric": "sql"},   # classify -> 'sql'
    )
    # generate_sql uses complete too; script the NL2SQL text via same needle.
    client._complete["the primary metric for"] = "SELECT * FROM gold_metrics"
    captured = {}
    router = QueryRouter(client=client,
                         sql_executor=lambda q: captured.setdefault("q", q) or [])
    out = _run(router.answer("Show the primary metric for my firms", team_filter="capital"))
    assert out["sql"] is not None
    assert "scoped" in captured["q"]        # team scope wrapper applied


# --- reviewer loop --------------------------------------------------------

def test_reviewer_override_promotes_to_gold():
    claim = Claim(field_name="the primary metric", canonical_metric="primary_ratio",
                  value="14.2", confidence=0.5, needs_review=True,
                  source_element_id="e1", entity_ref="FIRM123")
    dec = review("e1", "14.2", bbox=None, reviewer="alice",
                 override_value="15.0")
    apply_decision(claim, dec)
    assert claim.value == "15.0"
    assert claim.needs_review is False
    assert claim.citation_tier.value == "manual"          # overridden = human-typed
    assert claim.review_tier.value == "review_override"
    rows = DeltaLakeStorage(root="/tmp/pra_test").build_gold_metric_rows([claim])
    assert rows[0]["value"] == "15.0" and rows[0]["entity_ref"] == "FIRM123"
    assert decision_row(dec)["review_tier"] == "review_override"


# --- storage idempotency --------------------------------------------------

def test_gold_write_is_idempotent(tmp_path):
    claim = Claim(field_name="the primary metric", canonical_metric="primary_ratio",
                  value="14.2", confidence=0.9, source_element_id="e1",
                  entity_ref="FIRM123")
    store = DeltaLakeStorage(root=str(tmp_path))
    assert store.write_gold_metrics([claim]) == 1
    assert store.write_gold_metrics([claim]) == 0   # dedup on re-run
    assert len(store.read_jsonl("gold")) == 1


# --- gates ----------------------------------------------------------------

def test_gates_enforce_signoff_and_validation():
    with pytest.raises(PermissionError):
        gate_external_release({"external": True}, None, Stage.PRODUCTION)
    ok = gate_external_release(
        {"external": True},
        SignOff("carol", "Head of Supervision", "2026-01-15T09:00:00"),
        Stage.PRODUCTION)
    assert ok["released"] is True
    with pytest.raises(PermissionError):
        gate_production_deployment(None)
    with pytest.raises(ValueError):
        SignOff("x", "y", "not-a-date")   # ISO-8601 enforced


# --- conflict indexing to search -----------------------------------------

def test_conflicts_indexed_to_search():
    from docextract.custom.production import _index_conflicts
    claims = [
        Claim(field_name="the primary metric", canonical_metric="primary_ratio", value="14",
              unit="%", source_element_id="e1", entity_ref="FIRM123"),
        Claim(field_name="the primary metric", canonical_metric="primary_ratio", value="14",
              unit="bps", source_element_id="e2", entity_ref="FIRM123"),
    ]
    store = FakeSearchStore()
    n = _index_conflicts(claims, {"path": "/x/doc.pdf"}, store, "FIRM123")
    assert n == 1 and len(store.docs) == 1


# --- crop-and-zoom (figure preprocessor) ----------------------------------

from docextract.search import figure_preprocessor as _fp
from docextract.shared.schema import CitationTier, ReviewTier
from docextract.shared.llm_client import LLMResponse as _LLMResp


def _pil_available():
    try:
        import PIL  # noqa: F401
        return True
    except ImportError:
        return False


def _figure_el(desc="Bar chart of the primary metric over time"):
    return Element(element_id="f1", modality=Modality.FIGURE, content="",
                   source_document="/x/d.pdf", team="capital",
                   report_type="primary_report", entity_ref="FIRM123",
                   page=1, bbox=(0.1, 0.1, 0.5, 0.5), description=desc)


def test_is_chart_like():
    assert _fp.is_chart_like(_figure_el()) is True
    assert _fp.is_chart_like(_figure_el(desc="a company logo")) is False
    text = Element(element_id="t", modality=Modality.TEXT, content="x",
                   source_document="/x/d.pdf", entity_ref="F1")
    assert _fp.is_chart_like(text) is False


@pytest.mark.skipif(not _pil_available(), reason="Pillow not installed")
def test_crop_and_zoom_upscales_and_caps():
    from PIL import Image
    page = Image.new("RGB", (1000, 800), "white")
    crop = _fp.crop_and_zoom(page, (0.25, 0.25, 0.35, 0.35))
    assert min(crop.size) >= 1024           # small region upscaled
    assert max(crop.size) <= 2000           # capped
    full = _fp.crop_and_zoom(page, (0.0, 0.0, 1.0, 1.0))
    assert max(full.size) <= 2000


class _RecMM:
    """Records whether an image reached the model per call."""
    def __init__(self):
        self.calls = []

    async def extract(self, model_key, prompt, content, tool_schema, image_base64=None):
        self.calls.append(image_base64 is not None)
        return _LLMResp(text="", tool_arguments={"claims": [
            {"field_name": "the primary metric", "value": "14.2", "confidence": 0.95}]},
            input_tokens=10, output_tokens=5)


class _FakeProvider:
    def get_page_image(self, document_bytes, filename, page_number):
        from PIL import Image
        return Image.new("RGB", (1000, 800), "white")


@pytest.mark.skipif(not _pil_available(), reason="Pillow not installed")
def test_chart_figure_read_multimodally_and_tagged():
    rc = _RecMM()
    res = _run(Extractor(client=rc).extract(
        [_figure_el()], document_bytes=b"PDF", filename="doc.pdf",
        page_image_provider=_FakeProvider()))
    assert rc.calls and rc.calls[0] is True          # image reached the model
    assert res.cropped_figures == 1
    c = res.claims[0]
    assert c.citation_tier == CitationTier.LLM_ESTIMATED
    assert c.page == 1 and c.bbox == (0.1, 0.1, 0.5, 0.5)


def test_figure_degrades_to_text_only_without_provider():
    rc = _RecMM()
    res = _run(Extractor(client=rc).extract([_figure_el()]))   # no provider
    assert rc.calls and rc.calls[0] is False
    assert res.cropped_figures == 0 and len(res.claims) == 1


def test_provider_returning_none_degrades():
    class _NoneProv:
        def get_page_image(self, *a):
            return None
    rc = _RecMM()
    res = _run(Extractor(client=rc).extract(
        [_figure_el()], document_bytes=b"X", filename="doc.pdf",
        page_image_provider=_NoneProv()))
    assert rc.calls and rc.calls[0] is False and len(res.claims) == 1


# --- two-axis tiers + structural classification ---------------------------

def test_text_provenance_parsed_figure_llm_estimated():
    client = StubLLMClient(default_claims=[
        {"field_name": "the primary metric", "value": "14.2", "confidence": 0.95}])
    text = _element("has the primary metric")
    assert _run(Extractor(client=client).extract([text])).claims[0].citation_tier \
        == CitationTier.PARSED


def test_reviewer_override_sets_manual_and_override_tier():
    claim = Claim(field_name="the primary metric", canonical_metric="primary_ratio",
                  value="14.2", confidence=0.5, needs_review=True,
                  source_element_id="e1", entity_ref="F1")
    dec = review("e1", "14.2", bbox=None, reviewer="a", override_value="15.0")
    apply_decision(claim, dec)
    assert claim.citation_tier == CitationTier.MANUAL
    assert claim.review_tier == ReviewTier.HUMAN_OVERRIDE
    assert decision_row(dec)["review_tier"] == "review_override"


def test_reviewer_confirm_keeps_provenance():
    claim = Claim(field_name="coverage", value="120", source_element_id="e2",
                  entity_ref="F1")
    apply_decision(claim, review("e2", "120", bbox=None, reviewer="b"))
    assert claim.review_tier == ReviewTier.HUMAN_CONFIRMED
    assert claim.citation_tier == CitationTier.PARSED


def test_metric_classifier_measure_and_period():
    from docextract.shared.metric_normaliser import classify_metric
    cl = classify_metric("Primary ratio YoY change", "14.2%")
    assert cl.measure_type.value == "change"
    assert cl.period.value == "year_on_year"


# --- claim-quant matcher --------------------------------------------------

from docextract.custom.claim_quant_matcher import (
    looks_like_claim, match_claim, match_all)


def _cq_claims():
    return [
        Claim(field_name="Primary ratio", canonical_metric="primary_ratio",
              value="14.2", source_element_id="e1", entity_ref="F1"),
        Claim(field_name="coverage", canonical_metric="coverage_ratio",
              value="120", source_element_id="e2", entity_ref="F1"),
    ]


class _CQClient:
    def __init__(self, idx, conf):
        self._idx, self._conf = idx, conf

    async def complete(self, model, prompt, content):
        import json
        return _LLMResp(text=json.dumps(
            {"best_candidate_index": self._idx, "confidence": self._conf}))


def test_claim_gate():
    assert looks_like_claim("Primary metric strengthened this quarter")
    assert not looks_like_claim("Section 3: Capital Overview")


def test_claim_link_resolves_above_threshold():
    link = _run(match_claim("Primary metric strengthened", _cq_claims(), _CQClient(0, 0.9)))
    assert link.matched_metric == "primary_ratio" and link.confidence == 0.9


def test_claim_link_unlinked_below_threshold():
    link = _run(match_claim("Primary metric strengthened", _cq_claims(), _CQClient(0, 0.4)))
    assert link.matched_metric is None       # not force-matched


def test_claim_link_null_not_forced():
    link = _run(match_claim("vague statement", _cq_claims(), _CQClient(None, 0.0)))
    assert link.matched_metric is None


def test_claim_link_malformed_json_safe():
    class _Bad:
        async def complete(self, m, p, c):
            return _LLMResp(text="not json")
    link = _run(match_claim("the primary metric rose", _cq_claims(), _Bad()))
    assert link.matched_metric is None


def test_match_all_filters_non_claims():
    links = _run(match_all(["Primary metric strengthened", "A heading", "coverage fell"],
                      _cq_claims(), _CQClient(0, 0.9)))
    assert len(links) == 2


# --- per-metric review thresholds -----------------------------------------

def test_per_metric_review_thresholds(tmp_path):
    import os as _o
    _o.environ["DOCEXTRACT_THRESHOLD_REGISTRY_DIR"] = str(tmp_path)
    try:
        from docextract.shared.config import CONFIG
        from docextract.shared.threshold_registry import (
            MetricThreshold, register_threshold, invalidate_cache)
        invalidate_cache()
        th = CONFIG.thresholds
        assert th.review_threshold("primary_ratio") == th.human_review  # empty -> global
        register_threshold(MetricThreshold(metric="primary_ratio", threshold=0.95),
                           updated_by="governance")
        assert th.review_threshold("primary_ratio") == 0.95
        assert th.review_threshold("unknown_metric") == th.human_review
    finally:
        invalidate_cache()
        del _o.environ["DOCEXTRACT_THRESHOLD_REGISTRY_DIR"]


def test_per_metric_threshold_respects_global_floor(tmp_path):
    import os as _o
    _o.environ["DOCEXTRACT_THRESHOLD_REGISTRY_DIR"] = str(tmp_path)
    try:
        from docextract.shared.config import CONFIG
        from docextract.shared.threshold_registry import (
            MetricThreshold, register_threshold, invalidate_cache)
        invalidate_cache()
        th = CONFIG.thresholds  # global 0.60
        register_threshold(MetricThreshold(metric="primary_ratio", threshold=0.95),
                           updated_by="governance")
        assert th.review_threshold("primary_ratio") == 0.95
        assert th.review_threshold(None) == 0.60
        register_threshold(MetricThreshold(metric="risky", threshold=0.40),
                           updated_by="governance")
        assert th.review_threshold("risky") == 0.60   # floored up
    finally:
        invalidate_cache()
        del _o.environ["DOCEXTRACT_THRESHOLD_REGISTRY_DIR"]


# --- multi-team storage naming --------------------------------------------

from docextract.shared.config import CONFIG as _CFG


def test_storage_slug_and_naming():
    sp = _CFG.storage
    assert sp.slug("fx", "exchange_report") == "fx_exchange_report"
    assert sp.slug("FX Team", "Exchange Report") == "fx_team_exchange_report"
    assert sp.layer_table("gold", "fx", "exchange_report") == \
        "docextract.reporting.gold_fx_exchange_report"
    assert sp.fqtn("ops_audit_events") == "docextract.reporting.ops_audit_events"
    assert sp.layer_local("silver", "fx", "exchange_report") == \
        "silver__fx_exchange_report"


def test_partitioned_write_keeps_entity_ref_as_column(tmp_path):
    import json
    store = DeltaLakeStorage(root=str(tmp_path))
    c = Claim(field_name="the primary metric", canonical_metric="primary_ratio", value="14.2",
              confidence=0.9, source_element_id="e1", entity_ref="FRN123")
    assert store.write_gold_metrics([c], team="fx",
                                    report_type="exchange_report") == 1
    f = tmp_path / "gold__fx_exchange_report.jsonl"
    assert f.exists()
    row = json.loads(f.read_text().splitlines()[0])
    assert row["entity_ref"] == "FRN123"          # entity_ref is a column, not the name


def test_partitioned_write_separates_teams(tmp_path):
    store = DeltaLakeStorage(root=str(tmp_path))
    c = Claim(field_name="X", value="1", source_element_id="e", entity_ref="F")
    store.write_silver([c], team="finance", report_type="primary_report")
    store.write_silver([c], team="operations", report_type="coverage_report")
    assert (tmp_path / "silver__finance_primary_report.jsonl").exists()
    assert (tmp_path / "silver__operations_coverage_report.jsonl").exists()


# --- 3-level orchestration discovery --------------------------------------

def test_entity_ref_from_filename():
    from docextract.custom.run_ingestion_job import entity_ref_from_filename
    assert entity_ref_from_filename(
        "FRN123_AcmeBank_CapAdq_2026-06-30.pdf") == "FRN123"
    assert entity_ref_from_filename("not_a_convention.pdf") is None


def test_discovery_scans_team_report_entity_layout(tmp_path):
    from docextract.custom.run_ingestion_job import discover_submissions
    import os
    os.makedirs(tmp_path / "fx" / "exchange_report" / "ACME_CORP")
    os.makedirs(tmp_path / "capital" / "adequacy" / "BETA_LTD")
    (tmp_path / "fx" / "exchange_report" / "ACME_CORP" /
     "anything named FINAL.pdf").write_text("x")
    (tmp_path / "capital" / "adequacy" / "BETA_LTD" / "q1.pdf").write_text("x")
    (tmp_path / "capital" / "adequacy" / "BETA_LTD" / ".hidden").write_text("x")
    subs = discover_submissions(str(tmp_path))
    assert len(subs) == 2
    fx = [s for s in subs if s["team"] == "fx"][0]
    assert fx["report_type"] == "exchange_report"
    assert fx["entity_ref"] == "ACME_CORP"          # from the FOLDER, not filename
    assert all(s["entity_ref"] for s in subs)


def test_loose_file_without_entity_folder_is_not_collected(tmp_path):
    from docextract.custom.run_ingestion_job import discover_submissions
    import os
    os.makedirs(tmp_path / "fx" / "exchange_report" / "ACME_CORP")
    (tmp_path / "fx" / "exchange_report" / "ACME_CORP" / "good.pdf").write_text("x")
    (tmp_path / "fx" / "exchange_report" / "STRAY.pdf").write_text("x")
    subs = discover_submissions(str(tmp_path))
    names = {os.path.basename(s["path"]) for s in subs}
    assert "good.pdf" in names
    assert "STRAY.pdf" not in names                 # no stray


def test_entity_folder_name_normalised(tmp_path):
    from docextract.custom.run_ingestion_job import discover_submissions
    import os
    os.makedirs(tmp_path / "fx" / "exchange_report" / "  Acme Corp  ")
    (tmp_path / "fx" / "exchange_report" / "  Acme Corp  " / "d.pdf").write_text("x")
    subs = discover_submissions(str(tmp_path))
    assert subs[0]["entity_ref"] == "Acme Corp"


# --- stored prompt registry -----------------------------------------------

def test_prompt_registry_default_seed_then_override(tmp_path):
    import os
    os.environ["PRA_PROMPT_REGISTRY_DIR"] = str(tmp_path)
    try:
        from docextract.shared import prompt_registry as pr
        d = pr.get_prompt("fx", "exchange_report", "query")
        assert d.team == "default"
        pr.register_prompt(pr.PromptEntry(
            team="fx", report_type="exchange_report", kind="query",
            preamble="FX preamble",
            fewshots=[pr.FewShot("Q", "A")]), updated_by="alice")
        g = pr.get_prompt("fx", "exchange_report", "query")
        assert g.team == "fx" and "FX preamble" in g.preamble
        assert "BODY" in g.render("BODY") and "Q" in g.render("BODY")
    finally:
        del os.environ["PRA_PROMPT_REGISTRY_DIR"]


def test_prompt_registry_list_missing(tmp_path):
    import os
    os.environ["PRA_PROMPT_REGISTRY_DIR"] = str(tmp_path)
    try:
        from docextract.shared import prompt_registry as pr
        missing = pr.list_missing([("liquidity", "lcr")])
        assert ("liquidity", "lcr", "extraction") in missing
        assert ("liquidity", "lcr", "query") in missing
    finally:
        del os.environ["PRA_PROMPT_REGISTRY_DIR"]


# --- supervisor session: memory, reference-resolution, logging, feedback ---

from docextract.shared.supervisor_session import (
    SupervisorSession, load_query_log)


class _SessClient:
    """Scripted client for router+session tests."""
    async def complete(self, model, prompt, content):
        p = prompt.lower()
        if "self-contained" in p or "rewrite" in p:
            if "survive" in content.lower():
                return _LLMResp(text="Will firm A survive given its the primary metric of 8.2%?")
            return _LLMResp(text=content.split("Follow-up question:")[-1]
                            .split("Rewritten")[0].strip())
        if "classify" in p:
            return _LLMResp(text="sql")
        if "SELECT" in prompt:
            return _LLMResp(text="SELECT entity_ref, canonical_metric, value FROM gold_metrics")
        return _LLMResp(text="answer")


def _sess_router():
    from docextract.shared.llm_router import QueryRouter
    return QueryRouter(
        client=_SessClient(),
        sql_executor=lambda q: [{"entity_ref": "FIRM_A",
                                 "canonical_metric": "primary_ratio", "value": "8.2"}])


def test_first_turn_not_rewritten_and_recorded(tmp_path):
    import os
    os.environ["PRA_SUPERVISOR_LOG_DIR"] = str(tmp_path)
    try:
        sess = SupervisorSession(session_id="s1")
        r = _sess_router()
        a = _run(r.answer("What is the financial risk of firm A",
                     team_filter="capital", session=sess))
        assert a["resolved_question"] == a["raw_question"]   # no history yet
        assert a.get("turn_id") and sess.has_history()
        assert sess.recent_turns()[-1].key_figures            # figures captured
    finally:
        del os.environ["PRA_SUPERVISOR_LOG_DIR"]


def test_followup_reference_is_resolved(tmp_path):
    import os
    os.environ["PRA_SUPERVISOR_LOG_DIR"] = str(tmp_path)
    try:
        sess = SupervisorSession(session_id="s2")
        r = _sess_router()
        _run(r.answer("What is the financial risk of firm A", team_filter="capital", session=sess))
        a2 = _run(r.answer("Will firm A survive from the figure", team_filter="capital", session=sess))
        # the follow-up is rewritten into a self-contained question
        assert a2["resolved_question"] != a2["raw_question"]
        assert "8.2" in a2["resolved_question"] or "the primary metric" in a2["resolved_question"]
    finally:
        del os.environ["PRA_SUPERVISOR_LOG_DIR"]


def test_qa_persisted_to_audit_log(tmp_path):
    import os
    os.environ["PRA_SUPERVISOR_LOG_DIR"] = str(tmp_path)
    try:
        sess = SupervisorSession(session_id="s3")
        r = _sess_router()
        _run(r.answer("What is the financial risk of firm A", team_filter="capital", session=sess))
        log = load_query_log()
        assert any(x.get("log_kind") == "answer" for x in log)
        assert any(x.get("resolved_question") for x in log)
    finally:
        del os.environ["PRA_SUPERVISOR_LOG_DIR"]


def test_feedback_capture_and_validation(tmp_path):
    import os
    os.environ["PRA_SUPERVISOR_LOG_DIR"] = str(tmp_path)
    try:
        sess = SupervisorSession(session_id="s4")
        r = _sess_router()
        a = _run(r.answer("What is the financial risk of firm A", team_filter="capital", session=sess))
        t = sess.add_feedback(a["turn_id"], "down", note="missed leverage")
        assert t is not None and t.feedback_rating == "down"
        assert any(x.get("feedback_rating") == "down" for x in load_query_log())
        with pytest.raises(ValueError):
            sess.add_feedback(a["turn_id"], "maybe")
        # feedback on an expired/unknown turn still logs, returns None
        assert sess.add_feedback("nope", "up") is None
    finally:
        del os.environ["PRA_SUPERVISOR_LOG_DIR"]


def test_stateless_call_still_works():
    # No session passed -> single-shot, no rewrite, no crash.
    r = _sess_router()
    a = _run(r.answer("What is the financial risk of firm A", team_filter="capital"))
    assert a["resolved_question"] == a["raw_question"] and "route" in a


def test_supervisor_feedback_summary(tmp_path):
    import os
    os.environ["PRA_SUPERVISOR_LOG_DIR"] = str(tmp_path)
    try:
        sess = SupervisorSession(session_id="s5")
        r = _sess_router()
        a1 = _run(r.answer("risk of firm A", team_filter="capital", session=sess))
        a2 = _run(r.answer("risk of firm B", team_filter="capital", session=sess))
        sess.add_feedback(a1["turn_id"], "up")
        sess.add_feedback(a2["turn_id"], "down")
        from docextract.shared.monitoring import supervisor_feedback_summary
        s = supervisor_feedback_summary()
        assert s["questions"] == 2
        assert s["thumbs_up"] == 1 and s["thumbs_down"] == 1
    finally:
        del os.environ["PRA_SUPERVISOR_LOG_DIR"]


# --- footnote preservation + text/bbox indexing ---------------------------

from docextract.custom.preprocessor import (
    _split_paragraphs, _footnote_markers, _paragraphs_with_bbox)
from docextract.custom.production import _index_text_elements


def test_footnote_region_kept_whole():
    text = "Body one.\n\nBody two.\n\n(a) first note\ncontinues.\n(b) second.\n(c) third."
    paras = _split_paragraphs(text)
    # body paragraphs stay separate; the whole footnote block is ONE chunk
    assert paras[0] == "Body one." and paras[1] == "Body two."
    footnote_chunks = [p for p in paras if p.startswith("(a)")]
    assert len(footnote_chunks) == 1
    assert all(m in footnote_chunks[0] for m in ("(a)", "(b)", "(c)"))


def test_plain_text_split_unchanged():
    assert _split_paragraphs("one.\n\ntwo.\n\nthree.") == ["one.", "two.", "three."]


def test_footnote_markers_extracted():
    assert _footnote_markers("row\n(a) note\n(f) note") == ["a", "f"]
    assert _footnote_markers("no markers here") == []


class _FakePage:
    width = 600
    height = 800
    def extract_text(self):
        return "Hello world.\n\n(a) footnote here."
    def extract_words(self):
        return [
            {"x0": 10, "top": 20, "x1": 60, "bottom": 40, "text": "Hello"},
            {"x0": 65, "top": 20, "x1": 120, "bottom": 40, "text": "world."},
            {"x0": 10, "top": 700, "x1": 40, "bottom": 720, "text": "(a)"},
            {"x0": 45, "top": 700, "x1": 120, "bottom": 720, "text": "footnote"},
            {"x0": 125, "top": 700, "x1": 160, "bottom": 720, "text": "here."},
        ]


def test_text_elements_carry_normalised_bbox():
    pairs = _paragraphs_with_bbox(_FakePage())
    assert len(pairs) == 2
    _, body_bbox = pairs[0]
    _, fn_bbox = pairs[1]
    assert body_bbox is not None and all(0 <= v <= 1 for v in body_bbox)
    assert fn_bbox is not None and fn_bbox[1] > body_bbox[1]   # footnote lower


def test_bbox_degrades_without_words():
    class _NoWords(_FakePage):
        def extract_words(self):
            return []
    assert _paragraphs_with_bbox(_NoWords())[0][1] is None


class _FakeSearchStore:
    def __init__(self):
        self.docs = []
    def index_many(self, chunks):
        self.docs.extend(chunks)


def test_text_and_footnotes_reach_vector_store():
    tbl = Element(element_id="d-p7-tbl0", modality=Modality.TABLE,
                  content="SLR 5.0\n(a) minimum 3.0%\n(f) buffer",
                  source_document="/x/d.pdf", page=7, bbox=(0.1, 0.2, 0.9, 0.5),
                  entity_ref="E1")
    fn = Element(element_id="d-p7-t3", modality=Modality.TEXT,
                 content="(a) SLR minimum 3.0%.\n(f) buffer 2.0%.",
                 source_document="/x/d.pdf", page=7, bbox=(0.1, 0.55, 0.9, 0.7),
                 entity_ref="E1")
    empty = Element(element_id="d-p7-t9", modality=Modality.TEXT, content="  ",
                    source_document="/x/d.pdf", page=7, entity_ref="E1")
    store = _FakeSearchStore()
    n = _index_text_elements([tbl, fn, empty], {"path": "/x/d.pdf"}, store, "E1")
    assert n == 2 and len(store.docs) == 2            # empty skipped
    tbl_chunk = next(c for c in store.docs if "SLR 5.0" in c.content)
    fn_chunk = next(c for c in store.docs if "minimum 3.0%." in c.content)
    # page + bbox travel with the chunk
    assert tbl_chunk.metrics["page"] == 7
    assert tbl_chunk.metrics["bbox"] == [0.1, 0.2, 0.9, 0.5]
    # shared footnote markers link table <-> footnote
    assert set(tbl_chunk.metrics["footnote_markers"]) == {"a", "f"}
    assert set(fn_chunk.metrics["footnote_markers"]) == {"a", "f"}


def test_index_text_elements_empty_safe():
    store = _FakeSearchStore()
    empty = Element(element_id="e", modality=Modality.TEXT, content="",
                    source_document="/x", page=1, entity_ref="E1")
    assert _index_text_elements([empty], {"path": "/x"}, store, "E1") == 0


# --- pptx + email preprocessing -------------------------------------------

def _pptx_available():
    try:
        import pptx  # noqa
        return True
    except ImportError:
        return False


@pytest.mark.skipif(not _pptx_available(), reason="python-pptx not installed")
def test_pptx_captures_text_table_notes(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    from docextract.custom.preprocessor import preprocess
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "Summary"
    s1.placeholders[1].text = "Primary ratio improved."
    s1.notes_slide.notes_text_frame.text = "Note the trend."
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    tb = s2.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(1)).table
    tb.cell(0, 0).text = "Metric"; tb.cell(1, 1).text = "15.6"
    p = tmp_path / "deck.pptx"; prs.save(str(p))
    els = preprocess(str(p), "finance", "primary_report", entity_ref="E1")
    assert any("Summary" in e.content or "improved" in e.content for e in els)
    assert any(e.modality == Modality.TABLE and "15.6" in e.content for e in els)
    assert any("Note the trend" in e.content for e in els)      # speaker notes
    assert all(e.entity_ref == "E1" for e in els)
    assert all(e.page in (1, 2) for e in els if e.page)


def test_eml_captures_header_and_body_ignores_attachment(tmp_path):
    from email.message import EmailMessage
    from docextract.custom.preprocessor import preprocess
    msg = EmailMessage()
    msg["From"] = "a@x.com"; msg["Subject"] = "Q1 figures"
    msg["Date"] = "Tue, 1 Apr 2025 10:00:00 +0000"
    msg.set_content("Primary ratio is 15.6%.\n\nReview by Friday.")
    # add an attachment — it must be IGNORED, not unpacked into elements
    msg.add_attachment(b"PDFDATA", maintype="application",
                       subtype="pdf", filename="report.pdf")
    p = tmp_path / "mail.eml"; p.write_bytes(bytes(msg))
    els = preprocess(str(p), "finance", "correspondence", entity_ref="E2")
    assert any("Q1 figures" in e.content for e in els)          # header
    assert any("Primary ratio is 15.6%" in e.content for e in els)  # body
    assert all(e.modality == Modality.TEXT for e in els)
    # attachment content never appears as an element
    assert not any("PDFDATA" in e.content or "report.pdf" in e.content for e in els)


def test_msg_without_library_flagged_not_crashed(tmp_path):
    import importlib.util
    from docextract.custom.preprocessor import preprocess
    if importlib.util.find_spec("extract_msg") is not None:
        return  # library present; fallback path not exercised
    p = tmp_path / "x.msg"; p.write_bytes(b"not a real msg")
    els = preprocess(str(p), "finance", "correspondence", entity_ref="E3")
    assert len(els) == 1 and els[0].modality == Modality.IMAGE   # flagged empty


def test_new_formats_registered():
    from docextract.custom.preprocessor import HANDLERS
    for ext in (".pptx", ".eml", ".msg"):
        assert ext in HANDLERS


# --- automatic eval against independent ground truth ----------------------

import os as _os
from docextract.shared.eval_registry import (
    ExpectedValue, register_expected, get_expected, list_missing)
from docextract.custom.evaluator import (
    is_sampled, evaluate_document, maybe_evaluate)


def test_eval_registry_register_and_retrieve(tmp_path):
    _os.environ["DOCEXTRACT_EVAL_REGISTRY_DIR"] = str(tmp_path)
    try:
        register_expected(ExpectedValue(
            team="finance", report_type="primary_report", entity_ref="E1",
            metric="primary_ratio", expected_value="15.6", tolerance=0.05),
            updated_by="tester")
        exp = get_expected("finance", "primary_report", "E1", "primary_ratio")
        assert exp is not None and exp.matches("15.62")
        assert not exp.matches("16.0")
        # unregistered metric returns None (no invented answer)
        assert get_expected("finance", "primary_report", "E1", "nope") is None
    finally:
        del _os.environ["DOCEXTRACT_EVAL_REGISTRY_DIR"]


def test_eval_grades_against_registered_and_catches_wrong(tmp_path):
    _os.environ["DOCEXTRACT_EVAL_REGISTRY_DIR"] = str(tmp_path)
    try:
        register_expected(ExpectedValue(
            team="finance", report_type="primary_report", entity_ref="E1",
            metric="primary_ratio", expected_value="15.6", tolerance=0.05),
            updated_by="tester")
        good = [Claim(field_name="Primary ratio", canonical_metric="primary_ratio",
                      value="15.61", source_element_id="e", entity_ref="E1")]
        assert evaluate_document("d.pdf", good, "finance", "primary_report", "E1").accuracy == 1.0
        bad = [Claim(field_name="Primary ratio", canonical_metric="primary_ratio",
                     value="99.9", source_element_id="e", entity_ref="E1")]
        assert evaluate_document("d.pdf", bad, "finance", "primary_report", "E1").accuracy == 0.0
    finally:
        del _os.environ["DOCEXTRACT_EVAL_REGISTRY_DIR"]


def test_eval_never_counts_ungraded_as_correct(tmp_path):
    _os.environ["DOCEXTRACT_EVAL_REGISTRY_DIR"] = str(tmp_path)
    try:
        ung = [Claim(field_name="M", canonical_metric="unregistered",
                     value="1", source_element_id="e", entity_ref="E1")]
        res = evaluate_document("d.pdf", ung, "finance", "primary_report", "E1")
        assert res.graded == 0 and res.ungraded == 1 and res.accuracy is None
    finally:
        del _os.environ["DOCEXTRACT_EVAL_REGISTRY_DIR"]


def test_eval_uses_human_review_as_truth():
    hr = [
        Claim(field_name="x", canonical_metric="m1", value="1",
              source_element_id="e", entity_ref="E1",
              review_tier=ReviewTier.HUMAN_CONFIRMED),
        Claim(field_name="y", canonical_metric="m2", value="2",
              source_element_id="e", entity_ref="E1",
              review_tier=ReviewTier.HUMAN_OVERRIDE),
    ]
    res = evaluate_document("d.pdf", hr, "finance", "primary_report", "E1")
    assert res.human_reviewed == 2 and res.human_overrides == 1
    assert res.override_rate == 0.5


def test_eval_sampling_deterministic_and_fractional():
    ids = [f"doc{i}.pdf" for i in range(1000)]
    sampled = [d for d in ids if is_sampled(d, 10)]
    assert 70 < len(sampled) < 130                     # ~1/10
    assert is_sampled("doc5.pdf", 10) == is_sampled("doc5.pdf", 10)  # stable
    assert is_sampled("anything", 1)                   # N<=1 samples all


def test_list_missing_flags_unregistered(tmp_path):
    _os.environ["DOCEXTRACT_EVAL_REGISTRY_DIR"] = str(tmp_path)
    try:
        register_expected(ExpectedValue(
            team="finance", report_type="primary_report", entity_ref="E1",
            metric="primary_ratio", expected_value="15.6"), updated_by="t")
        miss = list_missing([("finance", "primary_report"),
                             ("ops", "coverage_report")])
        assert ("ops", "coverage_report") in miss
        assert ("finance", "primary_report") not in miss
    finally:
        del _os.environ["DOCEXTRACT_EVAL_REGISTRY_DIR"]


# --- dictionary as a registry (data, not code) ----------------------------

from docextract.shared.dictionary import canonical_metric, canonical_doctype
from docextract.shared.dictionary_registry import (
    TermMapping, register_term, get_mappings, list_registered)


def test_dictionary_seed_works_with_empty_registry(tmp_path):
    _os.environ["DOCEXTRACT_DICTIONARY_REGISTRY_DIR"] = str(tmp_path)
    try:
        assert canonical_metric("primary ratio") == ("primary_ratio", True)
        assert canonical_metric("Unknown Thing") == ("unknown_thing", False)
        assert canonical_doctype("primary report") == "primary_report"
    finally:
        del _os.environ["DOCEXTRACT_DICTIONARY_REGISTRY_DIR"]


def test_registered_term_maps_without_code_change(tmp_path):
    _os.environ["DOCEXTRACT_DICTIONARY_REGISTRY_DIR"] = str(tmp_path)
    try:
        register_term(TermMapping(kind="metric", synonym="Return on Equity",
                                  canonical="roe"), updated_by="t")
        assert canonical_metric("return on equity") == ("roe", True)
        assert canonical_metric("  RETURN  ON EQUITY ") == ("roe", True)  # normalised
    finally:
        del _os.environ["DOCEXTRACT_DICTIONARY_REGISTRY_DIR"]


def test_registry_overrides_seed_and_latest_wins(tmp_path):
    _os.environ["DOCEXTRACT_DICTIONARY_REGISTRY_DIR"] = str(tmp_path)
    try:
        register_term(TermMapping(kind="metric", synonym="coverage",
                                  canonical="v1"), updated_by="a")
        assert canonical_metric("coverage") == ("v1", True)   # overrides seed
        register_term(TermMapping(kind="metric", synonym="coverage",
                                  canonical="v2"), updated_by="b")
        assert canonical_metric("coverage") == ("v2", True)   # latest wins
    finally:
        del _os.environ["DOCEXTRACT_DICTIONARY_REGISTRY_DIR"]


def test_dictionary_kind_validation_and_separation(tmp_path):
    _os.environ["DOCEXTRACT_DICTIONARY_REGISTRY_DIR"] = str(tmp_path)
    try:
        register_term(TermMapping(kind="doctype", synonym="Annual Filing",
                                  canonical="annual_filing"), updated_by="t")
        assert canonical_doctype("annual filing") == "annual_filing"
        assert "annual filing" in list_registered("doctype")
        assert "annual filing" not in list_registered("metric")
        with pytest.raises(ValueError):
            register_term(TermMapping(kind="bogus", synonym="x", canonical="y"),
                          updated_by="t")
    finally:
        del _os.environ["DOCEXTRACT_DICTIONARY_REGISTRY_DIR"]


# --- default + registered query few-shots for the supervisor app ----------

from docextract.shared.prompt_registry import (
    get_prompt, register_prompt, PromptEntry, FewShot)


def test_default_query_prompt_has_fewshots():
    d = get_prompt("unregistered", "primary_report", "query")
    assert len(d.fewshots) >= 3                       # ships with examples
    assert any("human-verified" in f.output_text for f in d.fewshots)
    # no rename scars in the default text
    assert "a the pipeline" not in d.preamble
    assert "for review" not in d.preamble


def test_team_can_register_unlimited_query_fewshots(tmp_path):
    _os.environ["PRA_PROMPT_REGISTRY_DIR"] = str(tmp_path)
    try:
        register_prompt(PromptEntry(
            team="finance", report_type="primary_report", kind="query",
            preamble="Finance query guidance.",
            fewshots=[FewShot(input_text=f"Q{i}", output_text=f"A{i}")
                      for i in range(9)]),
            updated_by="t")
        f = get_prompt("finance", "primary_report", "query")
        assert len(f.fewshots) == 9                    # unbounded
        assert f.preamble.startswith("Finance")        # team's own, pulled per-team
        # a different team still gets the default
        assert get_prompt("ops", "coverage_report", "query").team == "default"
    finally:
        del _os.environ["PRA_PROMPT_REGISTRY_DIR"]


# --- content-domain tagging cascade (advisory) ----------------------------

def test_tag_cascade_heading_metric_llm_unknown(tmp_path):
    import os as _o, asyncio as _a
    _o.environ["DOCEXTRACT_TAG_REGISTRY_DIR"] = str(tmp_path)
    try:
        from docextract.shared.tag_registry import TagEntry, register_tag, UNKNOWN_TAG
        from docextract.custom.tagger import tag_element
        from docextract.shared.schema import Element, Claim, Modality

        for dom in ["risk", "financial", "tax"]:
            register_tag(TagEntry(kind="domain", value=dom), updated_by="gov")
        register_tag(TagEntry(kind="metric_domain", key="capital_ratio",
                              value="financial"), updated_by="gov")
        register_tag(TagEntry(kind="heading_domain", key="tax disclosures",
                              value="tax"), updated_by="gov")

        def E(content, desc=None):
            return Element(element_id="e", modality=Modality.TEXT, content=content,
                           source_document="d.pdf", description=desc)
        def C(m):
            return Claim(field_name=m, canonical_metric=m, value="1",
                         source_element_id="e")
        run = _a.run

        # rung 1: heading
        assert run(tag_element(E("x", desc="Part 3: Tax Disclosures"), [])).source == "heading"
        # rung 2: metric
        r2 = run(tag_element(E("capital"), [C("capital_ratio")]))
        assert r2.domain == "financial" and r2.source == "metric"
        # rung 4: nothing matches, no llm -> unknown (never crashes/guesses)
        assert run(tag_element(E("nothing"), [])).domain == UNKNOWN_TAG

        # rung 3 + safety: llm low-confidence falls to unknown
        class LLM:
            async def extract(self, m, p, c, t):
                class R: tool_arguments = {"domain": "risk", "confidence": 0.2}
                return R()
        assert run(tag_element(E("amb"), [], llm_client=LLM())).domain == UNKNOWN_TAG
    finally:
        del _o.environ["DOCEXTRACT_TAG_REGISTRY_DIR"]


def test_domain_tag_flows_to_gold_and_silver():
    from docextract.custom.storage import DeltaLakeStorage
    c = Claim(field_name="Capital ratio", canonical_metric="capital_ratio",
              value="15.6", source_element_id="e1", entity_ref="ACME",
              domain_tag="financial", tag_source="metric")
    gold = DeltaLakeStorage.build_gold_metric_rows([c])
    assert gold[0]["domain_tag"] == "financial"        # reaches curated Gold
    silver = c.model_dump(mode="json")
    assert silver["domain_tag"] == "financial"          # reaches Silver
    assert silver["tag_source"] == "metric"
    # untagged claim -> None, no crash
    c2 = Claim(field_name="x", canonical_metric="y", value="1", source_element_id="e2")
    assert DeltaLakeStorage.build_gold_metric_rows([c2])[0]["domain_tag"] is None


def test_databricks_search_result_parser():
    # the search store rewrite (Databricks AI Search) must return the same
    # {text, metadata, score} shape the callers expect, mapping columns by name
    from docextract.search.ai_search_store import _parse_results
    resp = {
        "manifest": {"columns": [
            {"name": "chunk_id"}, {"name": "content"}, {"name": "chunk_type"},
            {"name": "entity_ref"}, {"name": "source_document_id"},
            {"name": "content_hash"}, {"name": "score"}]},
        "result": {"data_array": [
            ["c1", "liquidity ratio 1.5", "table", "ACME", "doc1", "h1", 0.87]]},
    }
    out = _parse_results(resp)
    assert len(out) == 1
    assert out[0]["text"] == "liquidity ratio 1.5"
    assert out[0]["metadata"]["entity_ref"] == "ACME"
    assert out[0]["score"] == 0.87
    assert set(out[0].keys()) == {"text", "metadata", "score"}
    # malformed / empty -> [] not crash
    assert _parse_results({}) == []
