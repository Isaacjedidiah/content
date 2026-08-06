"""Format-agnostic preprocessing.

The only module that knows about file formats. Everything downstream works
on ``Element`` objects. Unknown formats fail loudly rather than silently
dropping a submission.

Real parsing uses pdfplumber (PDF text/tables), python-docx (Word),
openpyxl (Excel), python-pptx (PowerPoint), and the stdlib email package
(.eml; .msg needs the optional extract-msg library). Imports are lazy so the
module can be imported in environments lacking a given parser.

Email attachments are deliberately NOT unpacked here: teams place attachments
in the submission folder as their own files, so each flows through the
pipeline independently with its own provenance.

Scanned / image-only PDF pages have no text layer, so pdfplumber returns
nothing for them. Those pages fall back to ``ai_parse_document`` (the same
``databricks-ai-parse`` endpoint the native track uses) for OCR, keeping one
OCR path across both tracks. If the endpoint is unavailable or yields
nothing, the page keeps its flagged-empty element and goes to quarantine —
no silent data loss.

entity_ref resolution: the firm reference is threaded onto every Element from
the submission descriptor (resolved upstream from the volume path or a
sidecar manifest). This is what makes Gold attributable — previously nothing
populated it.
"""
from __future__ import annotations

import os
import re
from typing import Callable, Optional

from ..shared.schema import Element, Modality
from .ai_parse_client import AIParseClient


def preprocess_pdf(path: str, team: str, report_type: str,
                   entity_ref: Optional[str] = None,
                   ocr_client: Optional[AIParseClient] = None) -> list[Element]:
    import pdfplumber  # lazy

    # Instantiate the OCR fallback once per document. If it isn't configured
    # (e.g. local dev without a Databricks endpoint), ``available`` is False
    # and we skip straight to the flagged-empty behaviour.
    ocr = ocr_client if ocr_client is not None else AIParseClient()

    elements: list[Element] = []
    with pdfplumber.open(path) as pdf:
        for page_no, page in enumerate(pdf.pages, start=1):
            paragraphs = _paragraphs_with_bbox(page)   # (text, bbox) pairs
            tables = page.extract_tables() or []

            if not paragraphs and not tables:
                # No text layer: likely scanned. Try OCR via ai_parse_document;
                # fall back to a flagged-empty element (-> quarantine) if that
                # is unavailable or returns nothing.
                ocr_elements = _ocr_page(page, page_no, path, team,
                                         report_type, entity_ref, ocr)
                elements.extend(ocr_elements)
                continue

            for para, bbox in paragraphs:
                elements.append(_text_element(
                    path, page_no, len(elements), para, team, report_type,
                    entity_ref, bbox=bbox))
            for tbl_no, table in enumerate(tables):
                elements.append(_table_element(
                    path, page_no, tbl_no, _table_to_markdown(table), team,
                    report_type, entity_ref))
    return elements


def _ocr_page(page, page_no: int, path: str, team: str, report_type: str,
              entity_ref: Optional[str],
              ocr: AIParseClient) -> list[Element]:
    """OCR one text-less page, or return a single flagged-empty element.

    The flagged-empty element preserves the previous behaviour exactly, so a
    page we still can't read is never silently dropped — it goes to
    quarantine downstream.
    """
    flagged_empty = [Element(
        element_id=f"{os.path.basename(path)}-p{page_no}-empty",
        modality=Modality.IMAGE,
        content="",
        source_document=path,
        page=page_no,
        team=team,
        report_type=report_type,
        entity_ref=entity_ref,
    )]

    if not ocr.available:
        return flagged_empty

    try:
        image_bytes = _render_page_png(page)
        blocks = ocr.parse_page_image(image_bytes, mime="image/png")
    except Exception:
        # Any parse/transport failure: keep the page visible via quarantine
        # rather than losing it. (No silent fallback to a stub.)
        return flagged_empty

    if not blocks:
        return flagged_empty

    out: list[Element] = []
    for i, block in enumerate(blocks):
        if block.is_table:
            out.append(_table_element(
                path, page_no, i, block.content, team, report_type, entity_ref,
                suffix="ocr"))
        else:
            out.append(_text_element(
                path, page_no, i, block.content, team, report_type, entity_ref,
                suffix="ocr"))
    return out


def _render_page_png(page) -> bytes:
    """Render a pdfplumber page to PNG bytes for OCR.

    Uses pdfplumber's page.to_image(); the resolution is high enough for OCR
    without producing an oversized payload.
    """
    import io

    img = page.to_image(resolution=200)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _text_element(path: str, page_no: int, idx: int, content: str, team: str,
                  report_type: str, entity_ref: Optional[str],
                  suffix: str = "t",
                  bbox: Optional[tuple] = None) -> Element:
    return Element(
        element_id=f"{os.path.basename(path)}-p{page_no}-{suffix}{idx}",
        modality=Modality.TEXT,
        content=content,
        source_document=path,
        page=page_no,
        bbox=bbox,
        team=team,
        report_type=report_type,
        entity_ref=entity_ref,
    )


def _table_element(path: str, page_no: int, idx: int, content: str, team: str,
                   report_type: str, entity_ref: Optional[str],
                   suffix: str = "tbl") -> Element:
    return Element(
        element_id=f"{os.path.basename(path)}-p{page_no}-{suffix}{idx}",
        modality=Modality.TABLE,
        content=content,
        source_document=path,
        page=page_no,
        team=team,
        report_type=report_type,
        entity_ref=entity_ref,
    )


def preprocess_docx(path: str, team: str, report_type: str,
                    entity_ref: Optional[str] = None) -> list[Element]:
    import docx  # python-docx, lazy

    doc = docx.Document(path)
    elements: list[Element] = []
    for i, para in enumerate(doc.paragraphs):
        if para.text.strip():
            elements.append(Element(
                element_id=f"{os.path.basename(path)}-para{i}",
                modality=Modality.TEXT,
                content=para.text,
                source_document=path,
                team=team,
                report_type=report_type,
                entity_ref=entity_ref,
            ))
    for t, table in enumerate(doc.tables):
        rows = [[c.text for c in row.cells] for row in table.rows]
        elements.append(Element(
            element_id=f"{os.path.basename(path)}-tbl{t}",
            modality=Modality.TABLE,
            content=_rows_to_markdown(rows),
            source_document=path,
            team=team,
            report_type=report_type,
            entity_ref=entity_ref,
        ))
    return elements


def preprocess_xlsx(path: str, team: str, report_type: str,
                    entity_ref: Optional[str] = None) -> list[Element]:
    from openpyxl import load_workbook  # lazy

    wb = load_workbook(path, data_only=True)
    elements: list[Element] = []
    for sheet in wb.worksheets:
        rows = [[("" if c is None else str(c)) for c in row]
                for row in sheet.iter_rows(values_only=True)]
        if rows:
            elements.append(Element(
                element_id=f"{os.path.basename(path)}-{sheet.title}",
                modality=Modality.TABLE,
                content=_rows_to_markdown(rows),
                source_document=path,
                team=team,
                report_type=report_type,
                entity_ref=entity_ref,
            ))
    return elements


def preprocess_pptx(path: str, team: str, report_type: str,
                    entity_ref: Optional[str] = None) -> list[Element]:
    """Read slide text, tables and speaker notes from a .pptx.

    Slide TEXT is cheap (local, no model call). Slides also often carry
    meaning in images/diagrams — those are flagged as FIGURE elements with the
    slide as ``page`` so the downstream vision path can handle them, exactly as
    chart figures are handled in PDFs. Slide number is used as ``page`` so a
    retrieved chunk can be traced to its slide.
    """
    from pptx import Presentation  # python-pptx, lazy

    prs = Presentation(path)
    # slide dimensions (EMU) for normalising picture bboxes to 0-1
    sw = float(prs.slide_width or 1)
    sh = float(prs.slide_height or 1)

    elements: list[Element] = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        for shp_no, shape in enumerate(slide.shapes):
            # tables
            if shape.has_table:
                rows = [[cell.text for cell in row.cells]
                        for row in shape.table.rows]
                elements.append(Element(
                    element_id=f"{os.path.basename(path)}-s{slide_no}-tbl{shp_no}",
                    modality=Modality.TABLE,
                    content=_rows_to_markdown(rows),
                    source_document=path, page=slide_no,
                    team=team, report_type=report_type, entity_ref=entity_ref))
                continue
            # text frames
            if shape.has_text_frame and shape.text_frame.text.strip():
                elements.append(_text_element(
                    path, slide_no, shp_no, shape.text_frame.text, team,
                    report_type, entity_ref, suffix="s"))
                continue
            # pictures -> FIGURE with a normalised bbox, for the vision path
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    bbox = (float(shape.left) / sw, float(shape.top) / sh,
                            float(shape.left + shape.width) / sw,
                            float(shape.top + shape.height) / sh)
                except (TypeError, ValueError):
                    bbox = None
                elements.append(Element(
                    element_id=f"{os.path.basename(path)}-s{slide_no}-img{shp_no}",
                    modality=Modality.FIGURE,
                    content="",
                    source_document=path, page=slide_no, bbox=bbox,
                    team=team, report_type=report_type, entity_ref=entity_ref))

        # speaker notes as text (often carry the real narrative)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                elements.append(_text_element(
                    path, slide_no, 999, notes, team, report_type,
                    entity_ref, suffix="notes"))
    return elements


def preprocess_email(path: str, team: str, report_type: str,
                     entity_ref: Optional[str] = None) -> list[Element]:
    """Read an email body (.eml / .msg), headers as metadata.

    ATTACHMENTS ARE DELIBERATELY IGNORED — teams place attachments in the
    submission folder as their own files, so they flow through the pipeline on
    their own (and get their own entity_ref / provenance) rather than being
    silently unpacked here. Only the message body and its headers are read.

    .eml is parsed with the stdlib ``email`` package. .msg (Outlook) needs the
    optional ``extract-msg`` library; if it's absent, the file is flagged-empty
    (-> quarantine) rather than crashing, mirroring the OCR fallback posture.
    """
    ext = os.path.splitext(path)[1].lower()
    headers: dict = {}
    body = ""

    if ext == ".eml":
        import email
        from email import policy
        with open(path, "rb") as fh:
            msg = email.message_from_binary_file(fh, policy=policy.default)
        headers = {k: str(msg[k]) for k in ("from", "to", "subject", "date")
                   if msg[k] is not None}
        part = msg.get_body(preferencelist=("plain", "html"))
        body = part.get_content().strip() if part is not None else ""
    elif ext == ".msg":
        try:
            import extract_msg  # optional, lazy
        except ImportError:
            return [Element(
                element_id=f"{os.path.basename(path)}-empty",
                modality=Modality.IMAGE, content="",
                source_document=path, page=1, team=team,
                report_type=report_type, entity_ref=entity_ref)]
        m = extract_msg.Message(path)
        headers = {"from": m.sender or "", "to": m.to or "",
                   "subject": m.subject or "", "date": str(m.date or "")}
        body = (m.body or "").strip()

    elements: list[Element] = []
    # header block as its own text element (searchable context: who/when/subject)
    header_text = "\n".join(f"{k.title()}: {v}" for k, v in headers.items() if v)
    if header_text:
        elements.append(_text_element(
            path, 1, 0, header_text, team, report_type, entity_ref,
            suffix="hdr"))
    # body as paragraphs (footnote-aware split reused; emails rarely have
    # footnotes but the whole-block behaviour is harmless)
    for i, para in enumerate(_split_paragraphs(body), start=1):
        elements.append(_text_element(
            path, 1, i, para, team, report_type, entity_ref, suffix="body"))
    return elements


HANDLERS: dict[str, Callable[..., list[Element]]] = {
    ".pdf": preprocess_pdf,
    ".docx": preprocess_docx,
    ".xlsx": preprocess_xlsx,
    ".xlsm": preprocess_xlsx,
    ".pptx": preprocess_pptx,
    ".eml": preprocess_email,
    ".msg": preprocess_email,
}


def preprocess(path: str, team: str, report_type: str,
               entity_ref: Optional[str] = None,
               ocr_client: Optional[AIParseClient] = None) -> list[Element]:
    ext = os.path.splitext(path)[1].lower()
    if ext not in HANDLERS:
        raise ValueError(
            f"No preprocessor for {ext!r}. Add a handler or quarantine the file."
        )
    # Only the PDF handler performs OCR fallback; other formats have a real
    # text layer by construction, so the OCR client is not threaded to them.
    if ext == ".pdf":
        return preprocess_pdf(path, team, report_type, entity_ref,
                              ocr_client=ocr_client)
    return HANDLERS[ext](path, team, report_type, entity_ref)


# --- helpers ---------------------------------------------------------------

# A footnote marker at the start of a line: (a), (1), (i) etc. Used to detect
# the footnote region so its parts are kept together as one chunk rather than
# fragmented by blank-line splitting.
_FOOTNOTE_MARKER = re.compile(r"^\s*\(([a-z0-9]{1,3})\)\s", re.IGNORECASE)


def _footnote_markers(text: str) -> list[str]:
    """Return the footnote marker ids present at line-starts in ``text`` — e.g.
    a table containing '(a)' and '(f)' references. Used to link a table chunk
    to its footnote chunk even if they end up apart (see production indexing)."""
    out = []
    for line in text.splitlines():
        m = _FOOTNOTE_MARKER.match(line)
        if m:
            mid = m.group(1).lower()
            if mid not in out:
                out.append(mid)
    return out


def _split_paragraphs(text: str) -> list[str]:
    """Footnote-aware paragraph split (text-only; no coordinates).

    Splits the body on blank lines as before, BUT once the footnote region
    begins (first line starting with a ``(x)`` marker) everything from there to
    the end is kept together as a SINGLE paragraph, so a multi-part footnote
    block (a)(b)(c)... survives whole as one chunk instead of fragmenting.
    """
    lines = text.split("\n")
    footnote_start = None
    for i, line in enumerate(lines):
        if _FOOTNOTE_MARKER.match(line):
            footnote_start = i
            break

    if footnote_start is None:
        return [p.strip() for p in text.split("\n\n") if p.strip()]

    body = "\n".join(lines[:footnote_start])
    footnotes = "\n".join(lines[footnote_start:]).strip()
    paras = [p.strip() for p in body.split("\n\n") if p.strip()]
    if footnotes:
        paras.append(footnotes)          # whole footnote region = one chunk
    return paras


def _paragraphs_with_bbox(page) -> list[tuple[str, Optional[tuple]]]:
    """Build (paragraph_text, normalised_bbox) pairs from a pdfplumber page.

    Uses word coordinates so each paragraph carries its position on the page
    (normalised 0-1, matching how figure bboxes are stored). Falls back to
    ``(text, None)`` if word extraction is unavailable, so behaviour degrades
    to the previous text-only path rather than failing.
    """
    text = page.extract_text() or ""
    paragraphs = _split_paragraphs(text)
    if not paragraphs:
        return []

    try:
        words = page.extract_words() or []
        pw, ph = float(page.width), float(page.height)
    except Exception:
        return [(p, None) for p in paragraphs]

    if not words or not pw or not ph:
        return [(p, None) for p in paragraphs]

    # For each paragraph, consume its token count from the word stream and take
    # the bounding box of that span. Approximate (word order ~ reading order),
    # but good enough to locate a paragraph/footnote on the page.
    out: list[tuple[str, Optional[tuple]]] = []
    cursor = 0
    for para in paragraphs:
        para_tokens = para.split()
        if not para_tokens:
            out.append((para, None))
            continue
        span = words[cursor:cursor + len(para_tokens)]
        cursor += len(para_tokens)
        if not span:
            out.append((para, None))
            continue
        x0 = min(float(w["x0"]) for w in span) / pw
        top = min(float(w["top"]) for w in span) / ph
        x1 = max(float(w["x1"]) for w in span) / pw
        bottom = max(float(w["bottom"]) for w in span) / ph
        out.append((para, (x0, top, x1, bottom)))
    return out


def _table_to_markdown(table: list[list]) -> str:
    rows = [[("" if c is None else str(c)) for c in row] for row in table]
    return _rows_to_markdown(rows)


def _rows_to_markdown(rows: list[list[str]]) -> str:
    """Render rows as markdown, tolerating ragged rows.

    Real PDF tables from pdfplumber are frequently ragged; we pad short rows
    and never let a long row desync the column count.
    """
    if not rows:
        return ""
    width = max(len(r) for r in rows)

    def pad(r: list[str]) -> list[str]:
        return list(r) + [""] * (width - len(r))

    header = pad(rows[0])
    md = ["| " + " | ".join(header) + " |",
          "|" + "|".join(["---"] * width) + "|"]
    for row in rows[1:]:
        md.append("| " + " | ".join(pad(row)) + " |")
    return "\n".join(md)
